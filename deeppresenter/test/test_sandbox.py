"""
Merged Sandbox Testing Suite

This test suite includes the most comprehensive test for each tool category:
1. Matplotlib: Complex multi-subplot with Chinese text
2. Mermaid: Sequence diagram with Chinese text
3. MCP Tools: Complete workflow (create→edit→move)
4. Data Science: Full visualization pipeline (pandas+numpy+matplotlib+seaborn)
5. Image Processing: OpenCV advanced processing
6. Document Generation: python-pptx with multiple slides
7. System Tools: Ripgrep search with multiple files
8. Integration: Tools availability check

Each test represents the most challenging scenario for its category.
"""

import pytest
from pathlib import Path


class TestMatplotlibAdvanced:
    """Test matplotlib with the most complex scenario."""

    @pytest.mark.asyncio
    async def test_matplotlib_complex_chinese(
        self, agent_env, workspace, matplotlib_complex_script, tool_call_helper
    ):
        """Test complex matplotlib: 4 subplots + Chinese text + multiple chart types."""
        script_path = workspace / "test_complex.py"
        write_call = tool_call_helper(
            "write_file",
            {"path": str(script_path), "content": matplotlib_complex_script},
        )
        write_result = await agent_env.tool_execute(write_call)
        assert not write_result.is_error, f"Failed to write script: {write_result.text}"

        # Execute the script
        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {script_path}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"Script execution failed: {exec_result.text}"
        assert "SUCCESS" in exec_result.text, f"Script did not report success: {exec_result.text}"

        # Verify the output file exists
        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "complex_plot.png" in list_result.text, "Complex plot not found"


class TestMermaidAdvanced:
    """Test mermaid with the most complex scenario."""

    @pytest.mark.asyncio
    async def test_mermaid_sequence_chinese(
        self, agent_env, workspace, mermaid_sequence_diagram, tool_call_helper
    ):
        """Test mermaid sequence diagram with Chinese text and multiple participants."""
        # Write the sequence diagram
        diagram_path = workspace / "sequence.mmd"
        write_call = tool_call_helper(
            "write_file",
            {"path": str(diagram_path), "content": mermaid_sequence_diagram},
        )
        write_result = await agent_env.tool_execute(write_call)
        assert not write_result.is_error, f"Failed to write diagram: {write_result.text}"

        # Create Puppeteer config to allow running Chromium as root
        puppeteer_config_path = workspace / ".puppeteerrc.json"
        puppeteer_config = '{"args":["--no-sandbox","--disable-setuid-sandbox"]}'
        config_call = tool_call_helper(
            "write_file",
            {"path": str(puppeteer_config_path), "content": puppeteer_config},
        )
        config_result = await agent_env.tool_execute(config_call)
        assert not config_result.is_error, f"Failed to write Puppeteer config: {config_result.text}"

        # Generate PNG using mmdc with Puppeteer config
        output_path = workspace / "sequence_diagram.png"
        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"mmdc -i {diagram_path} -o {output_path} -p {puppeteer_config_path}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"mmdc execution failed: {exec_result.text}"
        assert "Error:" not in exec_result.text, f"mmdc reported error: {exec_result.text}"

        # Verify the output file exists
        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "sequence_diagram.png" in list_result.text, "Sequence diagram not found"


class TestMCPToolsWorkflow:
    """Test MCP tools with a complete workflow."""

    @pytest.mark.asyncio
    async def test_complete_workflow_with_chinese(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test complete workflow: nested dirs + write + edit + move + read (with Chinese)."""
        # Step 1: Create nested directory structure
        create_dir = tool_call_helper(
            "create_directory",
            {"path": str(workspace / "项目" / "源代码")},
        )
        await agent_env.tool_execute(create_dir)

        create_backup = tool_call_helper(
            "create_directory",
            {"path": str(workspace / "项目" / "备份")},
        )
        await agent_env.tool_execute(create_backup)

        # Step 2: Write initial file with Chinese content
        write_call = tool_call_helper(
            "write_file",
            {
                "path": str(workspace / "项目" / "源代码" / "应用.py"),
                "content": "# 版本 1.0\n# 作者：测试\nprint('你好，世界！')\n",
            },
        )
        write_result = await agent_env.tool_execute(write_call)
        assert not write_result.is_error

        # Step 3: Edit the file (multiple replacements)
        edit_call = tool_call_helper(
            "edit_file",
            {
                "file_path": str(workspace / "项目" / "源代码" / "应用.py"),
                "old_string": "版本 1.0",
                "new_string": "版本 2.0",
            },
        )
        edit_result = await agent_env.tool_execute(edit_call)
        assert not edit_result.is_error

        # Step 4: Append more content
        append_call = tool_call_helper(
            "write_file",
            {
                "path": str(workspace / "项目" / "源代码" / "应用.py"),
                "content": "\n# 新增功能\nprint('新版本发布')\n",
                "mode": "append",
            },
        )
        await agent_env.tool_execute(append_call)

        # Step 5: Move to backup with new name
        move_call = tool_call_helper(
            "move_file",
            {
                "source": str(workspace / "项目" / "源代码" / "应用.py"),
                "destination": str(workspace / "项目" / "备份" / "应用_v2.0.py"),
            },
        )
        move_result = await agent_env.tool_execute(move_call)
        assert not move_result.is_error

        # Step 6: Verify final state with list_directory
        list_backup = tool_call_helper(
            "list_directory",
            {"path": str(workspace / "项目" / "备份")},
        )
        list_result = await agent_env.tool_execute(list_backup)
        assert "应用_v2.0.py" in list_result.text

        # Step 7: Verify content with read_file
        read_call = tool_call_helper(
            "read_file",
            {"path": str(workspace / "项目" / "备份" / "应用_v2.0.py")},
        )
        read_result = await agent_env.tool_execute(read_call)
        assert "版本 2.0" in read_result.text
        assert "新增功能" in read_result.text
        assert "你好，世界！" in read_result.text


class TestDataSciencePipeline:
    """Test data science tools with complete pipeline."""

    @pytest.mark.asyncio
    async def test_full_data_visualization_pipeline(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test complete pipeline: numpy+pandas+matplotlib+seaborn with Chinese labels."""
        script_content = """
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns

# Configure Chinese fonts
plt.rcParams['font.sans-serif'] = ['Noto Sans CJK SC', 'WenQuanYi Zen Hei', 'DejaVu Sans']
plt.rcParams['axes.unicode_minus'] = False

# Generate sample data
np.random.seed(42)
data = pd.DataFrame({
    '日期': pd.date_range('2024-01-01', periods=100),
    '数值': np.random.randn(100).cumsum(),
    '类别': np.random.choice(['产品A', '产品B', '产品C'], 100),
    '销量': np.random.randint(10, 100, 100)
})

# Save raw data to CSV
data.to_csv('原始数据.csv', index=False, encoding='utf-8')

# Create complex visualizations
fig, axes = plt.subplots(2, 2, figsize=(14, 12))

# Plot 1: Time series line plot
axes[0, 0].plot(data['日期'], data['数值'], linewidth=2, color='steelblue')
axes[0, 0].set_title('时间序列趋势图', fontsize=14, fontweight='bold')
axes[0, 0].set_xlabel('日期', fontsize=12)
axes[0, 0].set_ylabel('累计数值', fontsize=12)
axes[0, 0].grid(True, alpha=0.3)

# Plot 2: Distribution histogram with seaborn
axes[0, 1].hist(data['数值'], bins=20, edgecolor='black', alpha=0.7, color='coral')
axes[0, 1].set_title('数值分布直方图', fontsize=14, fontweight='bold')
axes[0, 1].set_xlabel('数值范围', fontsize=12)
axes[0, 1].set_ylabel('频次', fontsize=12)

# Plot 3: Seaborn box plot by category
sns.boxplot(data=data, x='类别', y='销量', ax=axes[1, 0], palette='Set2')
axes[1, 0].set_title('各产品销量箱线图', fontsize=14, fontweight='bold')
axes[1, 0].set_xlabel('产品类别', fontsize=12)
axes[1, 0].set_ylabel('销量', fontsize=12)

# Plot 4: Seaborn scatter plot with regression
for cat in data['类别'].unique():
    cat_data = data[data['类别'] == cat]
    axes[1, 1].scatter(cat_data.index, cat_data['数值'], label=cat, alpha=0.6, s=50)
axes[1, 1].set_title('散点分布图（按类别）', fontsize=14, fontweight='bold')
axes[1, 1].set_xlabel('时间索引', fontsize=12)
axes[1, 1].set_ylabel('数值', fontsize=12)
axes[1, 1].legend(fontsize=10)
axes[1, 1].grid(True, alpha=0.3)

plt.tight_layout()
plt.savefig('数据可视化报告.png', dpi=100, bbox_inches='tight')

# Create detailed summary report
summary = f'''数据分析与可视化报告
{'='*50}

数据集信息:
- 总记录数: {len(data)}
- 日期范围: {data["日期"].min().strftime("%Y-%m-%d")} 至 {data["日期"].max().strftime("%Y-%m-%d")}
- 数值范围: {data["数值"].min():.2f} 至 {data["数值"].max():.2f}
- 平均数值: {data["数值"].mean():.2f}
- 标准差: {data["数值"].std():.2f}

各类别统计:
{data.groupby("类别")["销量"].describe().to_string()}

生成文件:
- 原始数据.csv: 原始数据导出
- 数据可视化报告.png: 综合可视化图表

使用工具:
- NumPy {np.__version__}
- Pandas {pd.__version__}
- Matplotlib {matplotlib.__version__}
- Seaborn {sns.__version__}
'''

with open('分析报告.txt', 'w', encoding='utf-8') as f:
    f.write(summary)

print('Pipeline completed successfully')
print(f'Files created: 原始数据.csv, 数据可视化报告.png, 分析报告.txt')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "pipeline.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'pipeline.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"pipeline failed: {exec_result.text}"
        assert "Pipeline completed successfully" in exec_result.text

        # Verify all files created
        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "原始数据.csv" in list_result.text
        assert "数据可视化报告.png" in list_result.text
        assert "分析报告.txt" in list_result.text


class TestImageProcessingAdvanced:
    """Test advanced image processing with OpenCV."""

    @pytest.mark.asyncio
    async def test_opencv_complex_processing(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test OpenCV: create image + draw multiple shapes + add text + apply filters."""
        script_content = """
import cv2
import numpy as np

# Create base image
img = np.zeros((400, 600, 3), dtype=np.uint8)
img[:, :] = (255, 255, 255)  # White background

# Draw complex shapes
cv2.rectangle(img, (50, 50), (550, 350), (0, 0, 255), 3)  # Red rectangle
cv2.circle(img, (300, 200), 100, (0, 255, 0), -1)  # Green filled circle
cv2.ellipse(img, (300, 200), (150, 80), 45, 0, 360, (255, 0, 0), 2)  # Blue ellipse
cv2.line(img, (50, 50), (550, 350), (255, 0, 255), 2)  # Magenta line
cv2.line(img, (550, 50), (50, 350), (255, 255, 0), 2)  # Cyan line

# Add Chinese text
font = cv2.FONT_HERSHEY_SIMPLEX
cv2.putText(img, 'OpenCV Advanced', (180, 80), font, 1.2, (0, 0, 0), 2, cv2.LINE_AA)
cv2.putText(img, 'Test Image', (220, 320), font, 1, (0, 0, 0), 2, cv2.LINE_AA)

# Draw polygon
pts = np.array([[100, 150], [150, 100], [250, 120], [200, 180]], np.int32)
pts = pts.reshape((-1, 1, 2))
cv2.polylines(img, [pts], True, (0, 128, 255), 2)

# Save original
cv2.imwrite('opencv_original.png', img)

# Apply Gaussian blur
blurred = cv2.GaussianBlur(img, (15, 15), 0)
cv2.imwrite('opencv_blurred.png', blurred)

# Apply edge detection
gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
edges = cv2.Canny(gray, 50, 150)
cv2.imwrite('opencv_edges.png', edges)

print(f'opencv_original.png created ({img.shape[1]}x{img.shape[0]})')
print(f'opencv_blurred.png created')
print(f'opencv_edges.png created')
print(f'OpenCV version: {cv2.__version__}')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_opencv.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_opencv.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"opencv test failed: {exec_result.text}"
        assert "opencv_original.png created" in exec_result.text
        assert "opencv_blurred.png created" in exec_result.text
        assert "opencv_edges.png created" in exec_result.text

        # Verify all three images created
        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "opencv_original.png" in list_result.text
        assert "opencv_blurred.png" in list_result.text
        assert "opencv_edges.png" in list_result.text


class TestDocumentGeneration:
    """Test document generation with python-pptx."""

    @pytest.mark.asyncio
    async def test_python_pptx_comprehensive(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test python-pptx: multiple slides + text + shapes + Chinese content."""
        script_content = """
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title slide with Chinese
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]
title.text = "Python-PPTX 综合测试"
subtitle.text = "自动生成演示文稿\\n沙箱环境测试"

# Slide 2: Content slide with bullets
content_slide = prs.slides.add_slide(prs.slide_layouts[1])
title = content_slide.shapes.title
content = content_slide.placeholders[1]
title.text = "主要功能特性"
tf = content.text_frame
tf.text = "功能 1: 基础文本支持"
p = tf.add_paragraph()
p.text = "功能 2: 多级列表"
p.level = 1
p = tf.add_paragraph()
p.text = "功能 3: 中文字体渲染"
p.level = 1
p = tf.add_paragraph()
p.text = "功能 4: 自定义样式"

# Slide 3: Blank slide with shapes
blank_slide = prs.slides.add_slide(prs.slide_layouts[6])
shapes = blank_slide.shapes

# Add title
title_box = shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "图形与形状演示"
title_frame.paragraphs[0].font.size = Pt(32)
title_frame.paragraphs[0].font.bold = True

# Add rectangle
rect = shapes.add_shape(
    1,  # Rectangle
    Inches(1), Inches(2), Inches(3), Inches(2)
)
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(0, 128, 255)
rect_text = rect.text_frame
rect_text.text = "蓝色矩形"

# Add circle (oval)
circle = shapes.add_shape(
    9,  # Oval
    Inches(5), Inches(2), Inches(2), Inches(2)
)
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(255, 128, 0)
circle_text = circle.text_frame
circle_text.text = "圆形"

# Slide 4: Summary
summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
title = summary_slide.shapes.title
content = summary_slide.placeholders[1]
title.text = "总结"
tf = content.text_frame
tf.text = "✓ 成功创建多页演示文稿"
p = tf.add_paragraph()
p.text = "✓ 支持中文内容"
p = tf.add_paragraph()
p.text = "✓ 包含图形和形状"
p = tf.add_paragraph()
p.text = "✓ 自定义样式和颜色"

# Save presentation
prs.save('综合演示.pptx')
print(f'综合演示.pptx created ({len(prs.slides)} slides)')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_pptx.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_pptx.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"python-pptx test failed: {exec_result.text}"
        assert "综合演示.pptx created" in exec_result.text
        assert "4 slides" in exec_result.text

        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "综合演示.pptx" in list_result.text


class TestSystemToolsAdvanced:
    """Test advanced system tools usage."""

    @pytest.mark.asyncio
    async def test_ripgrep_complex_search(self, agent_env, workspace, tool_call_helper):
        """Test ripgrep: search across multiple files with patterns."""
        # Create multiple Python files with various patterns
        files_content = {
            "模块1.py": """
# 配置文件
API_KEY = "test_key_123"
DEBUG = True
VERSION = "1.0.0"

def 初始化():
    print("初始化模块1")
    return API_KEY
""",
            "模块2.py": """
# 工具函数
import requests

API_KEY = "test_key_456"
BASE_URL = "https://api.example.com"

def 获取数据(api_key=None):
    if api_key is None:
        api_key = API_KEY
    print(f"使用密钥: {api_key}")
""",
            "配置.py": """
# 全局配置
DATABASE_URL = "postgresql://localhost/db"
API_KEY = "prod_key_789"
SECRET_KEY = "secret_value"

配置字典 = {
    "调试模式": False,
    "日志级别": "INFO"
}
""",
        }

        # Write all files
        for filename, content in files_content.items():
            write_call = tool_call_helper(
                "write_file",
                {"path": str(workspace / filename), "content": content},
            )
            await agent_env.tool_execute(write_call)

        # Test 1: Search for API_KEY pattern
        search_call1 = tool_call_helper(
            "execute_command",
            {"command": f"cd {workspace} && rg 'API_KEY' --color never"},
        )
        result1 = await agent_env.tool_execute(search_call1)
        assert not result1.is_error, f"ripgrep failed: {result1.text}"
        assert "API_KEY" in result1.text
        # Should find at least 3 occurrences across files
        assert result1.text.count("API_KEY") >= 3

        # Test 2: Search for Chinese patterns
        search_call2 = tool_call_helper(
            "execute_command",
            {"command": f"cd {workspace} && rg '初始化|配置' --color never"},
        )
        result2 = await agent_env.tool_execute(search_call2)
        assert not result2.is_error
        assert "初始化" in result2.text or "配置" in result2.text

        # Test 3: Count matches
        search_call3 = tool_call_helper(
            "execute_command",
            {"command": f"cd {workspace} && rg 'def ' --count --color never"},
        )
        result3 = await agent_env.tool_execute(search_call3)
        assert not result3.is_error
        # Should find function definitions


class TestIntegration:
    """Integration tests for sandbox environment."""

    @pytest.mark.asyncio
    async def test_all_tools_available(self, agent_env):
        """Verify all required sandbox tools are available."""
        required_tools = {
            "MCP Tools": [
                "execute_command",
                "write_file",
                "read_file",
                "list_directory",
                "create_directory",
                "move_file",
                "edit_file",
            ],
        }

        available_tools = list(agent_env._tools_dict.keys())

        for category, tools in required_tools.items():
            for tool in tools:
                assert tool in available_tools, f"Required tool '{tool}' not available in {category}"

        print(f"\n✓ All {sum(len(v) for v in required_tools.values())} required tools available")

    @pytest.mark.asyncio
    async def test_environment_capabilities(self, agent_env, workspace, tool_call_helper):
        """Test that all key libraries and tools are working."""
        test_script = """
import sys
print("=== Python Environment ===")
print(f"Python: {sys.version.split()[0]}")

# Test data science imports
import numpy as np
import pandas as pd
import matplotlib
import seaborn as sns
import plotly
print(f"NumPy: {np.__version__}")
print(f"Pandas: {pd.__version__}")
print(f"Matplotlib: {matplotlib.__version__}")
print(f"Seaborn: {sns.__version__}")
print(f"Plotly: {plotly.__version__}")

# Test image processing imports
import cv2
from PIL import Image
print(f"OpenCV: {cv2.__version__}")
print(f"Pillow: {Image.__version__}")

# Test document generation
from pptx import Presentation
print(f"python-pptx: Available")

print("\\n✓ All libraries loaded successfully")
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_env.py"), "content": test_script},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_env.py'}"},
        )
        result = await agent_env.tool_execute(exec_call)
        assert not result.is_error, f"Environment test failed: {result.text}"
        assert "All libraries loaded successfully" in result.text

        # Test system tools - all should be available in sandbox
        system_tools = [
            ("git --version", "git version"),
            ("mmdc --version", ""),  # mmdc might not output to stdout
            ("rg --version", "ripgrep"),
            ("convert --version", "ImageMagick"),
        ]

        for cmd, expected in system_tools:
            tool_call = tool_call_helper("execute_command", {"command": cmd})
            result = await agent_env.tool_execute(tool_call)
            # Assert that critical tools are available
            tool_name = cmd.split()[0]
            assert not (result.is_error and "not found" in result.text), \
                f"Critical system tool '{tool_name}' not available in sandbox"

            # Verify expected output if specified
            if expected and not result.is_error:
                assert expected in result.text, \
                    f"Tool '{tool_name}' found but output unexpected: {result.text[:100]}"
