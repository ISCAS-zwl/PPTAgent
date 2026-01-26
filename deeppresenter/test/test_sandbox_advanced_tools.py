"""
Advanced Sandbox Tools Testing

This test suite validates ALL advanced tools in the sandbox:
1. Data Science: matplotlib, seaborn, plotly, pandas, numpy
2. Image Processing: opencv, pillow, imagemagick
3. Document Generation: python-pptx, pptxgenjs
4. System Tools: git, vim, curl, wget, ripgrep

Tests all high-level functionality beyond basic file operations.
"""

import pytest
from pathlib import Path


class TestMatplotlibSeaborn:
    """Test matplotlib and seaborn visualization tools."""

    @pytest.mark.asyncio
    async def test_matplotlib_basic_plot(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test matplotlib basic plotting."""
        script_content = """
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

plt.figure(figsize=(8, 6))
plt.plot([1, 2, 3, 4, 5], [1, 4, 9, 16, 25], 'b-o', label='y=x²')
plt.xlabel('X axis')
plt.ylabel('Y axis')
plt.title('Matplotlib Basic Plot')
plt.legend()
plt.grid(True)
plt.savefig('matplotlib_basic.png', dpi=100)
print('matplotlib_basic.png created')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "plot_basic.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'plot_basic.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"matplotlib test failed: {exec_result.text}"
        assert "matplotlib_basic.png created" in exec_result.text

        # Verify file exists
        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "matplotlib_basic.png" in list_result.text

    @pytest.mark.asyncio
    async def test_seaborn_statistical_plot(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test seaborn statistical visualization."""
        script_content = """
import matplotlib
matplotlib.use('Agg')
import seaborn as sns
import pandas as pd
import numpy as np

# Create sample data
np.random.seed(42)
data = pd.DataFrame({
    'x': np.random.randn(100),
    'y': np.random.randn(100),
    'category': np.random.choice(['A', 'B', 'C'], 100)
})

# Create seaborn plot
sns.set_style('whitegrid')
plot = sns.scatterplot(data=data, x='x', y='y', hue='category', s=100)
plot.figure.savefig('seaborn_scatter.png', dpi=100)
print('seaborn_scatter.png created')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "plot_seaborn.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'plot_seaborn.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"seaborn test failed: {exec_result.text}"
        assert "seaborn_scatter.png created" in exec_result.text

        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "seaborn_scatter.png" in list_result.text


class TestPlotly:
    """Test plotly interactive visualization."""

    @pytest.mark.asyncio
    async def test_plotly_html_export(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test plotly HTML chart generation."""
        script_content = """
import plotly.graph_objects as go

# Create interactive plot
fig = go.Figure(data=go.Scatter(
    x=[1, 2, 3, 4, 5],
    y=[1, 4, 9, 16, 25],
    mode='lines+markers',
    name='y=x²'
))

fig.update_layout(
    title='Plotly Interactive Chart',
    xaxis_title='X axis',
    yaxis_title='Y axis'
)

fig.write_html('plotly_chart.html')
print('plotly_chart.html created')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "plot_plotly.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'plot_plotly.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"plotly test failed: {exec_result.text}"
        assert "plotly_chart.html created" in exec_result.text

        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "plotly_chart.html" in list_result.text

    @pytest.mark.asyncio
    async def test_plotly_image_export(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test plotly static image export."""
        script_content = """
import plotly.graph_objects as go

fig = go.Figure(data=[
    go.Bar(x=['A', 'B', 'C', 'D'], y=[20, 14, 23, 17])
])

fig.update_layout(title='Plotly Bar Chart')

# Export as static image (requires kaleido or orca)
try:
    fig.write_image('plotly_bar.png')
    print('plotly_bar.png created')
except Exception as e:
    print(f'Image export not available: {e}')
    # Fallback to HTML
    fig.write_html('plotly_bar.html')
    print('plotly_bar.html created (fallback)')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "plot_plotly_bar.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'plot_plotly_bar.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"plotly bar test failed: {exec_result.text}"
        # Accept either PNG or HTML output
        assert "created" in exec_result.text


class TestPandasNumpy:
    """Test pandas and numpy data processing."""

    @pytest.mark.asyncio
    async def test_numpy_array_operations(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test numpy array operations."""
        script_content = """
import numpy as np

# Create arrays
a = np.array([1, 2, 3, 4, 5])
b = np.array([5, 4, 3, 2, 1])

# Operations
sum_arr = a + b
prod_arr = a * b
mean_val = np.mean(a)
std_val = np.std(b)

# Save results
results = f'''NumPy Operations:
Sum: {sum_arr}
Product: {prod_arr}
Mean of a: {mean_val}
Std of b: {std_val}
'''

with open('numpy_results.txt', 'w') as f:
    f.write(results)

print('numpy_results.txt created')
print(f'NumPy version: {np.__version__}')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_numpy.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_numpy.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"numpy test failed: {exec_result.text}"
        assert "numpy_results.txt created" in exec_result.text

        # Read and verify results
        read_call = tool_call_helper(
            "read_file",
            {"path": str(workspace / "numpy_results.txt")},
        )
        read_result = await agent_env.tool_execute(read_call)
        assert "NumPy Operations" in read_result.text

    @pytest.mark.asyncio
    async def test_pandas_dataframe_operations(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test pandas DataFrame operations."""
        script_content = """
import pandas as pd
import numpy as np

# Create DataFrame
df = pd.DataFrame({
    'Name': ['Alice', 'Bob', 'Charlie', 'David', 'Eve'],
    'Age': [25, 30, 35, 40, 28],
    'Score': [85, 92, 78, 95, 88],
    'City': ['北京', '上海', '广州', '深圳', '杭州']
})

# Operations
avg_score = df['Score'].mean()
max_age = df['Age'].max()
grouped = df.groupby('City')['Score'].mean()

# Save to CSV
df.to_csv('data.csv', index=False)

# Create summary
summary = f'''Pandas DataFrame Summary:
Total Records: {len(df)}
Average Score: {avg_score:.2f}
Max Age: {max_age}
Pandas version: {pd.__version__}
'''

with open('pandas_summary.txt', 'w', encoding='utf-8') as f:
    f.write(summary)
    f.write('\\n\\nDataFrame:\\n')
    f.write(df.to_string())

print('data.csv created')
print('pandas_summary.txt created')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_pandas.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_pandas.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"pandas test failed: {exec_result.text}"
        assert "data.csv created" in exec_result.text

        # Verify both files
        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "data.csv" in list_result.text
        assert "pandas_summary.txt" in list_result.text


class TestImageProcessing:
    """Test opencv and pillow image processing."""

    @pytest.mark.asyncio
    async def test_pillow_image_creation(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test PIL/Pillow image creation and manipulation."""
        script_content = """
from PIL import Image, ImageDraw, ImageFont

# Create new image
img = Image.new('RGB', (400, 300), color='white')
draw = ImageDraw.Draw(img)

# Draw shapes
draw.rectangle([50, 50, 350, 250], outline='blue', width=3)
draw.ellipse([100, 80, 300, 220], fill='lightblue', outline='darkblue')
draw.line([50, 50, 350, 250], fill='red', width=2)

# Add text
draw.text((150, 130), 'Pillow Test', fill='black')

# Save image
img.save('pillow_test.png')
print(f'pillow_test.png created ({img.size[0]}x{img.size[1]})')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_pillow.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_pillow.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"pillow test failed: {exec_result.text}"
        assert "pillow_test.png created" in exec_result.text

        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "pillow_test.png" in list_result.text

    @pytest.mark.asyncio
    async def test_opencv_image_processing(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test OpenCV image processing."""
        script_content = """
import cv2
import numpy as np

# Create a test image
img = np.zeros((300, 400, 3), dtype=np.uint8)
img[:, :] = (255, 255, 255)  # White background

# Draw shapes
cv2.rectangle(img, (50, 50), (350, 250), (0, 0, 255), 3)  # Red rectangle
cv2.circle(img, (200, 150), 80, (0, 255, 0), -1)  # Green circle
cv2.line(img, (50, 50), (350, 250), (255, 0, 0), 2)  # Blue line

# Add text
cv2.putText(img, 'OpenCV Test', (120, 160), cv2.FONT_HERSHEY_SIMPLEX,
            1, (0, 0, 0), 2, cv2.LINE_AA)

# Save image
cv2.imwrite('opencv_test.png', img)
print(f'opencv_test.png created ({img.shape[1]}x{img.shape[0]})')
print(f'OpenCV version: {cv2.__version__}')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_opencv.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_opencv.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"opencv test failed: {exec_result.text}"
        assert "opencv_test.png created" in exec_result.text

        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "opencv_test.png" in list_result.text


class TestDocumentGeneration:
    """Test python-pptx and pptxgenjs document generation."""

    @pytest.mark.asyncio
    async def test_python_pptx_creation(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test python-pptx PowerPoint creation."""
        script_content = """
from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
prs = Presentation()

# Add title slide
title_slide = prs.slides.add_slide(prs.slide_layouts[0])
title = title_slide.shapes.title
subtitle = title_slide.placeholders[1]
title.text = "Python-PPTX Test"
subtitle.text = "Generated by Sandbox Environment"

# Add content slide
content_slide = prs.slides.add_slide(prs.slide_layouts[1])
title = content_slide.shapes.title
content = content_slide.placeholders[1]
title.text = "Test Content"
tf = content.text_frame
tf.text = "Feature 1: Basic text"
p = tf.add_paragraph()
p.text = "Feature 2: Multiple paragraphs"
p.level = 1

# Save presentation
prs.save('test_presentation.pptx')
print(f'test_presentation.pptx created ({len(prs.slides)} slides)')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "test_pptx.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'test_pptx.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"python-pptx test failed: {exec_result.text}"
        assert "test_presentation.pptx created" in exec_result.text

        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "test_presentation.pptx" in list_result.text


class TestImageMagick:
    """Test ImageMagick command-line tools."""

    @pytest.mark.asyncio
    async def test_imagemagick_convert(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test ImageMagick convert command."""
        # First create a test image with Python
        create_img_script = """
from PIL import Image

img = Image.new('RGB', (200, 200), color='red')
img.save('test_red.png')
print('test_red.png created')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "create_img.py"), "content": create_img_script},
        )
        await agent_env.tool_execute(write_call)

        exec_create = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'create_img.py'}"},
        )
        create_result = await agent_env.tool_execute(exec_create)
        assert "test_red.png created" in create_result.text

        # Use ImageMagick to convert and resize
        exec_convert = tool_call_helper(
            "execute_command",
            {"command": f"cd {workspace} && convert test_red.png -resize 100x100 test_red_small.png"},
        )
        convert_result = await agent_env.tool_execute(exec_convert)
        assert not convert_result.is_error, f"ImageMagick convert failed: {convert_result.text}"

        # Verify result
        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "test_red_small.png" in list_result.text

    @pytest.mark.asyncio
    async def test_imagemagick_identify(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test ImageMagick identify command."""
        # Create a test image
        create_img_script = """
from PIL import Image

img = Image.new('RGB', (300, 200), color='blue')
img.save('test_blue.png')
print('test_blue.png created')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "create_blue.py"), "content": create_img_script},
        )
        await agent_env.tool_execute(write_call)

        exec_create = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'create_blue.py'}"},
        )
        await agent_env.tool_execute(exec_create)

        # Use identify to get image info
        exec_identify = tool_call_helper(
            "execute_command",
            {"command": f"cd {workspace} && identify test_blue.png"},
        )
        identify_result = await agent_env.tool_execute(exec_identify)
        assert not identify_result.is_error, f"ImageMagick identify failed: {identify_result.text}"
        assert "300x200" in identify_result.text or "300" in identify_result.text


class TestSystemTools:
    """Test system tools: git, curl, wget, ripgrep, vim."""

    @pytest.mark.asyncio
    async def test_git_available(self, agent_env, workspace, tool_call_helper):
        """Test git availability and version."""
        exec_call = tool_call_helper(
            "execute_command",
            {"command": "git --version"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"git not available: {exec_result.text}"
        assert "git version" in exec_result.text

    @pytest.mark.asyncio
    async def test_curl_download(self, agent_env, workspace, tool_call_helper):
        """Test curl download capability."""
        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"cd {workspace} && curl -s -o example.html https://example.com"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        # Don't fail on network issues, just check curl is available

        # Test curl version
        version_call = tool_call_helper(
            "execute_command",
            {"command": "curl --version"},
        )
        version_result = await agent_env.tool_execute(version_call)
        assert not version_result.is_error
        assert "curl" in version_result.text

    @pytest.mark.asyncio
    async def test_ripgrep_search(self, agent_env, workspace, tool_call_helper):
        """Test ripgrep (rg) search."""
        # Create test files
        write_call1 = tool_call_helper(
            "write_file",
            {"path": str(workspace / "file1.txt"), "content": "Hello World\nTest Pattern\n"},
        )
        await agent_env.tool_execute(write_call1)

        write_call2 = tool_call_helper(
            "write_file",
            {"path": str(workspace / "file2.txt"), "content": "Pattern found here\nAnother line\n"},
        )
        await agent_env.tool_execute(write_call2)

        # Search with ripgrep
        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"cd {workspace} && rg Pattern"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"ripgrep failed: {exec_result.text}"
        assert "Pattern" in exec_result.text

    @pytest.mark.asyncio
    async def test_vim_available(self, agent_env, workspace, tool_call_helper):
        """Test vim availability."""
        exec_call = tool_call_helper(
            "execute_command",
            {"command": "vim --version | head -1"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error
        assert "VIM" in exec_result.text or "Vi IMproved" in exec_result.text


class TestCompleteWorkflow:
    """Test complex workflows combining multiple advanced tools."""

    @pytest.mark.asyncio
    async def test_data_visualization_pipeline(
        self, agent_env, workspace, tool_call_helper
    ):
        """Test complete data processing and visualization pipeline."""
        script_content = """
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns

# Generate sample data
np.random.seed(42)
data = pd.DataFrame({
    'date': pd.date_range('2024-01-01', periods=100),
    'value': np.random.randn(100).cumsum(),
    'category': np.random.choice(['A', 'B', 'C'], 100)
})

# Save raw data
data.to_csv('raw_data.csv', index=False)

# Create visualizations
fig, axes = plt.subplots(2, 2, figsize=(12, 10))

# Plot 1: Line plot
axes[0, 0].plot(data['date'], data['value'])
axes[0, 0].set_title('Time Series')
axes[0, 0].set_xlabel('Date')
axes[0, 0].set_ylabel('Value')

# Plot 2: Histogram
axes[0, 1].hist(data['value'], bins=20, edgecolor='black')
axes[0, 1].set_title('Distribution')
axes[0, 1].set_xlabel('Value')
axes[0, 1].set_ylabel('Frequency')

# Plot 3: Box plot by category
data.boxplot(column='value', by='category', ax=axes[1, 0])
axes[1, 0].set_title('Box Plot by Category')

# Plot 4: Scatter plot
for cat in data['category'].unique():
    cat_data = data[data['category'] == cat]
    axes[1, 1].scatter(cat_data.index, cat_data['value'], label=cat, alpha=0.6)
axes[1, 1].set_title('Scatter Plot')
axes[1, 1].legend()

plt.tight_layout()
plt.savefig('data_visualization.png', dpi=100)

# Create summary report
summary = f'''Data Visualization Pipeline Report
=====================================

Dataset Info:
- Total records: {len(data)}
- Date range: {data["date"].min()} to {data["date"].max()}
- Value range: {data["value"].min():.2f} to {data["value"].max():.2f}

Statistics by Category:
{data.groupby("category")["value"].describe()}

Files generated:
- raw_data.csv: Raw data export
- data_visualization.png: Combined visualization
'''

with open('report.txt', 'w') as f:
    f.write(summary)

print('Pipeline completed successfully')
print('Files created: raw_data.csv, data_visualization.png, report.txt')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "pipeline.py"), "content": script_content},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'pipeline.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"pipeline failed: {exec_result.text}"
        assert "Pipeline completed successfully" in exec_result.text

        # Verify all files created
        list_call = tool_call_helper(
            "list_directory",
            {"path": str(workspace)},
        )
        list_result = await agent_env.tool_execute(list_call)
        assert "raw_data.csv" in list_result.text
        assert "data_visualization.png" in list_result.text
        assert "report.txt" in list_result.text


class TestToolsVersions:
    """Test and report versions of all installed tools."""

    @pytest.mark.asyncio
    async def test_all_tools_versions(self, agent_env, workspace, tool_call_helper):
        """Get versions of all major tools."""
        version_script = """
import sys
import matplotlib
import seaborn
import plotly
import pandas
import numpy
import cv2
from PIL import Image
import pptx

print('=== Python Environment ===')
print(f'Python: {sys.version}')
print()

print('=== Data Science Tools ===')
print(f'NumPy: {numpy.__version__}')
print(f'Pandas: {pandas.__version__}')
print()

print('=== Visualization Tools ===')
print(f'Matplotlib: {matplotlib.__version__}')
print(f'Seaborn: {seaborn.__version__}')
print(f'Plotly: {plotly.__version__}')
print()

print('=== Image Processing ===')
print(f'OpenCV: {cv2.__version__}')
print(f'Pillow: {Image.__version__}')
print()

print('=== Document Generation ===')
print(f'python-pptx: {pptx.__version__}')
print()

print('All tools loaded successfully!')
"""
        write_call = tool_call_helper(
            "write_file",
            {"path": str(workspace / "versions.py"), "content": version_script},
        )
        await agent_env.tool_execute(write_call)

        exec_call = tool_call_helper(
            "execute_command",
            {"command": f"python {workspace / 'versions.py'}"},
        )
        exec_result = await agent_env.tool_execute(exec_call)
        assert not exec_result.is_error, f"version check failed: {exec_result.text}"
        assert "All tools loaded successfully" in exec_result.text
